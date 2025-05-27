import json
from pptx import Presentation

template_path = "template.pptx"
output_layout_json = "layout.json"

prs = Presentation(template_path)
layout_list = []
for idx, slide_layout in enumerate(prs.slide_layouts):
    layout_info = {
        "layout_index": idx,
        "layout_name": slide_layout.name,
        "placeholders": []
    }
    # Each layout has placeholder shapes
    for placeholder in slide_layout.placeholders:
        layout_info["placeholders"].append({
            "idx": placeholder.placeholder_format.idx,
            "name": placeholder.name,
            "type": placeholder.placeholder_format.type  # integer code for placeholder type
        })
    layout_list.append(layout_info)

with open(output_layout_json, 'w', encoding='utf-8') as f:
    json.dump({"layouts": layout_list}, f, indent=2, ensure_ascii=False)

print(f"Extracted layout saved to {output_layout_json}")

print("── Available slide layouts in template.pptx ──")
for idx, layout in enumerate(prs.slide_layouts):
    print(f"{idx:2d}: ‹{layout.name}›")
print("──────────────────────────────────────────────")