import json
from pptx import Presentation

TEMPLATE_PATH    = "template.pptx"
SLIDES_JSON_PATH = "slides.json"
OUTPUT_PATH      = "output.pptx"

with open(SLIDES_JSON_PATH, 'r', encoding='utf-8') as f:
    slides_def = json.load(f).get('slides', [])

prs = Presentation(TEMPLATE_PATH)

sldIdLst = prs.slides._sldIdLst
for sldId in list(sldIdLst):
    prs.part.drop_rel(sldId.rId)
    sldIdLst.remove(sldId)

layout_map = {layout.name: layout for layout in prs.slide_layouts}

for slide_def in slides_def:
    layout_name = slide_def.get('layout_name')
    if layout_name not in layout_map:
        raise ValueError(f"Unknown layout_name: '{layout_name}'. Valid names: {list(layout_map.keys())}")
    slide_layout = layout_map[layout_name]
    slide = prs.slides.add_slide(slide_layout)

    placeholders_dict = slide_def.get('placeholders', {})
    for idx_str, content in placeholders_dict.items():
        try:
            ph_idx = int(idx_str)
        except ValueError:
            print(f"⚠️ Invalid placeholder key (not an integer): '{idx_str}'")
            continue

        shape = next(
            (sh for sh in slide.placeholders if getattr(sh.placeholder_format, "idx", None) == ph_idx),
            None
        )
        if not shape:
            print(f"⚠️ Placeholder idx={ph_idx} not found on layout '{layout_name}'")
            continue

        tf = shape.text_frame
        if isinstance(content, list):
            tf.clear()
            for i, line in enumerate(content):
                if i == 0:
                    tf.text = line
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.level = 0
        else:
            tf.text = str(content)

prs.save(OUTPUT_PATH)
print(f"Target generated: {OUTPUT_PATH}")
